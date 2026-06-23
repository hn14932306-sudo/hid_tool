# -*- coding: utf-8 -*-
"""產生兩份圖文並茂的投影片：FAE 客戶版 與 完整功能版。
用 matplotlib 畫 UI 示意圖/流程圖，再用 python-pptx 組成 .pptx。"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, FancyArrowPatch, PathPatch
from matplotlib.path import Path
from PIL import Image

plt.rcParams["font.sans-serif"] = ["Microsoft JhengHei", "Microsoft YaHei", "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slides")
IMG = os.path.join(OUT, "img")
SHOTS = os.path.join(OUT, "shots")   # 使用者把實機截圖放這裡（檔名見下方清單）
os.makedirs(IMG, exist_ok=True)
os.makedirs(SHOTS, exist_ok=True)

# ---- 色系 ----
PRIMARY = "#22313f"   # 深藍灰（標題列）
ACCENT  = "#e74c3c"   # 紅（強調，呼應 heatmap icon）
ACCENT2 = "#3498db"   # 藍
GREEN   = "#2ecc71"
LIGHT   = "#f4f7f9"
PANEL   = "#ffffff"
BORDER  = "#cfd8dc"
TEXT    = "#22313f"
MUTED   = "#7f8c8d"
SLOT = ["#e74c3c", "#3498db", "#2ecc71", "#f39c12", "#9b59b6",
        "#1abc9c", "#e67e22", "#34495e", "#e91e63", "#00bcd4"]


def _ax(w=10.0, h=6.4):
    fig, ax = plt.subplots(figsize=(w, h), dpi=200)
    ax.set_xlim(0, 16); ax.set_ylim(0, 10); ax.axis("off")
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    return fig, ax


def _save(fig, name):
    p = os.path.join(IMG, name)
    fig.savefig(p, facecolor="white", dpi=200)
    plt.close(fig)
    return p


def rbox(ax, x, y, w, h, fc=PANEL, ec=BORDER, lw=1.2, rad=0.12, **kw):
    ax.add_patch(FancyBboxPatch((x, y), w, h,
                 boxstyle=f"round,pad=0.0,rounding_size={rad}",
                 fc=fc, ec=ec, lw=lw, **kw))


def window(ax, tabs, active, title="RE024 Touch Inspector"):
    rbox(ax, 0.3, 0.4, 15.4, 9.2, fc="white", ec=BORDER, lw=1.5, rad=0.18)
    ax.add_patch(Rectangle((0.3, 8.85), 15.4, 0.75, fc=PRIMARY, ec="none"))
    # 圓角上緣補丁
    rbox(ax, 0.3, 8.7, 15.4, 0.9, fc=PRIMARY, ec="none", rad=0.18)
    ax.text(0.65, 9.1, title, color="white", fontsize=12.5, va="center", fontweight="bold")
    ax.text(15.4, 9.1, "●", color=ACCENT, fontsize=11, va="center", ha="right")
    x = 0.55
    for i, t in enumerate(tabs):
        w = 2.05
        fc = ACCENT if i == active else "#eceff1"
        tc = "white" if i == active else "#455a64"
        rbox(ax, x, 7.95, w, 0.62, fc=fc, ec="none", rad=0.08)
        ax.text(x + w / 2, 8.26, t, color=tc, fontsize=10.5, ha="center", va="center",
                fontweight="bold" if i == active else "normal")
        x += w + 0.16


def arrow(ax, x1, y1, x2, y2, c=ACCENT2, lw=2.4):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                 mutation_scale=18, color=c, lw=lw))


def tri(ax, x, y, s=0.2, c=TEXT, left=False):
    pts = [(x + s, y + s), (x + s, y - s), (x - s, y)] if left \
        else [(x - s, y + s), (x - s, y - s), (x + s, y)]
    ax.add_patch(plt.Polygon(pts, closed=True, fc=c, ec="none"))


def check(ax, x, y, c=GREEN):
    ax.plot([x - 0.13, x - 0.02, x + 0.17], [y - 0.02, y - 0.15, y + 0.16],
            color=c, lw=2.6, solid_capstyle="round")


def cross(ax, x, y, c=ACCENT):
    ax.plot([x - 0.13, x + 0.13], [y - 0.13, y + 0.13], color=c, lw=2.6, solid_capstyle="round")
    ax.plot([x - 0.13, x + 0.13], [y + 0.13, y - 0.13], color=c, lw=2.6, solid_capstyle="round")


# ---------------------------------------------------------------- icon
def fig_icon():
    src = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "assets", "RE024_icon_heatmap.ico")
    p = os.path.join(IMG, "icon.png")
    try:
        im = Image.open(src)
        if hasattr(im, "size"):
            im = im.convert("RGBA")
        im.save(p)
    except Exception:
        p = None
    return p


# ---------------------------------------------------------------- window + monitor split
def fig_window():
    fig, ax = _ax()
    window(ax, ["監聽 Monitor", "發送 Send", "壓測 Stress", "回放 Replay"], 0)
    # 左：Monitor Data 表格
    rbox(ax, 0.6, 0.7, 7.1, 7.0, fc=LIGHT, ec=BORDER)
    ax.text(0.85, 7.35, "Monitor Data", fontsize=10.5, color=TEXT, fontweight="bold")
    cols = ["ID", "X", "Y", "CID", "Status", "Conf"]
    cx = [0.9, 2.0, 3.0, 4.0, 5.0, 6.7]
    ax.add_patch(Rectangle((0.75, 6.55), 6.8, 0.45, fc=PRIMARY, ec="none"))
    for c, t in zip(cx, cols):
        ax.text(c, 6.77, t, fontsize=8.5, color="white", va="center")
    rows = [["1", "487", "1303", "1", "Tip", "1"],
            ["2", "512", "1290", "2", "Tip", "1"],
            ["3", "—", "—", "—", "InRange", "0"],
            ["4", "933", "770", "1", "Tip Eraser", "1"]]
    for r, row in enumerate(rows):
        yy = 6.05 - r * 0.5
        if r % 2 == 0:
            ax.add_patch(Rectangle((0.75, yy - 0.2), 6.8, 0.45, fc="#e8eef2", ec="none"))
        for c, t in zip(cx, row):
            ax.text(c, yy, t, fontsize=8.2, color=TEXT, va="center")
    ax.text(0.9, 3.7, "RAW / FF01 欄位 · 全裝置或單一裝置", fontsize=8.6, color=MUTED)
    # 右：Touch Canvas（各佔一半）
    rbox(ax, 8.0, 0.7, 7.3, 7.0, fc="#fbfcfd", ec=BORDER)
    ax.text(8.25, 7.35, "Touch Canvas（各佔一半）", fontsize=10.5, color=TEXT, fontweight="bold")
    # 兩格裝置
    for gi, (gx, gy) in enumerate([(8.25, 4.1), (11.8, 4.1)]):
        rbox(ax, gx, gy, 3.25, 3.1, fc="white", ec="#dde6ea", rad=0.06)
        ax.text(gx + 0.15, gy + 2.85, f"Device {gi+1}", fontsize=8, color=SLOT[gi])
    # 觸控點 + 軌跡
    import numpy as np
    t = np.linspace(0, 1, 24)
    ax.plot(8.6 + 2.2 * t, 5.0 + 0.8 * np.sin(t * 6), color=SLOT[0], lw=2)
    for i, (px, py, col) in enumerate([(9.0, 5.6, SLOT[0]), (10.0, 4.9, SLOT[1]),
                                       (12.4, 6.2, SLOT[2]), (13.6, 5.0, SLOT[3])]):
        ax.add_patch(Circle((px, py), 0.16, fc=col, ec="white", lw=1.5))
        ax.text(px, py + 0.32, str(i + 1), fontsize=7.5, ha="center", color=col)
    # confidence=0 放大
    ax.add_patch(Circle((12.9, 4.5), 0.42, fc=SLOT[4], ec="black", lw=2))
    ax.text(12.9, 3.85, "Confidence=0\n放大", fontsize=7, ha="center", color=TEXT)
    return _save(fig, "window.png")


# ---------------------------------------------------------------- descriptor tree
def fig_descriptor():
    fig, ax = _ax(10, 6.4)
    rbox(ax, 0.4, 0.5, 15.2, 9.0, fc="white", ec=BORDER, lw=1.5)
    ax.text(0.8, 9.0, "Report Descriptor 欄位樹", fontsize=13, fontweight="bold", color=TEXT)
    nodes = [
        (0, "Usage Page: Digitizer (0x0D)", TEXT, True),
        (1, "Collection: TouchScreen", ACCENT2, True),
        (2, "ContactCount   (0x54)", MUTED, False),
        (2, "Collection: Finger ×10", ACCENT2, True),
        (3, "X / Y          (0x30/0x31)", TEXT, False),
        (3, "TipSwitch      (0x42)", TEXT, False),
        (3, "ContactID      (0x51)", TEXT, False),
        (3, "Confidence     (0x47)", TEXT, False),
        (1, "Collection: Pen", GREEN, True),
        (2, "X / Y · TipPressure (0x30)", TEXT, False),
        (2, "InRange · Invert · Eraser", TEXT, False),
        (2, "XTilt / YTilt / Azimuth", TEXT, False),
    ]
    y = 8.2
    for depth, text, col, bold in nodes:
        x = 1.0 + depth * 1.4
        if bold:
            ax.add_patch(Rectangle((x - 0.25, y - 0.02), 0.18, 0.18, fc=col, ec="none"))
        else:
            ax.add_patch(Circle((x - 0.16, y + 0.07), 0.05, fc=MUTED, ec="none"))
        ax.text(x + 0.1, y + 0.07, text, fontsize=10.2,
                color=col, va="center", fontweight="bold" if bold else "normal",
                family="Consolas")
        y -= 0.62
    ax.text(0.8, 0.75, "解析 HID Report Descriptor → 樹狀顯示每個欄位的 Usage / 位元位置 / 範圍",
            fontsize=9, color=MUTED)
    return _save(fig, "descriptor.png")


# ---------------------------------------------------------------- DigiInfo dual canvas
def fig_digiinfo():
    import numpy as np
    fig, ax = _ax()
    rbox(ax, 0.3, 0.4, 15.4, 9.2, fc="white", ec=BORDER, lw=1.5, rad=0.18)
    ax.text(0.7, 9.15, "回放 · DigiInfo（手 / 筆 雙畫布）", fontsize=12.5, fontweight="bold", color=TEXT)
    # 左 手(Touch)
    ax.text(0.9, 8.45, "手 (Touch)  0–2880 × 0–1800", fontsize=9.5, color=ACCENT2, fontweight="bold")
    rbox(ax, 0.8, 2.4, 6.9, 5.9, fc="#fbfcfd", ec=BORDER)
    th = np.linspace(0, 2.4, 60)
    for k in range(3):
        ax.plot(2.0 + 1.4 * np.cos(th + k) + k * 1.3, 5.0 + 1.6 * np.sin(th + k),
                color=SLOT[k], lw=2)
    ax.add_patch(Circle((4.6, 5.6), 0.16, fc=SLOT[0], ec="white", lw=1.5))
    # 右 筆(Pen)
    ax.text(8.3, 8.45, "筆 (Pen)  0–23040 × 0–14400", fontsize=9.5, color=GREEN, fontweight="bold")
    rbox(ax, 8.2, 2.4, 6.9, 5.9, fc="#fbfcfd", ec=BORDER)
    tp = np.linspace(0, 1, 80)
    ax.plot(9.2 + 5.0 * tp, 5.2 + 1.8 * np.sin(tp * 7) * (1 - tp), color=SLOT[4], lw=2.4)
    ax.add_patch(Circle((14.2, 5.2), 0.16, fc=SLOT[4], ec="white", lw=1.5))
    ax.text(9.4, 2.7, "只在 tip / eraser / 有壓力 時連線", fontsize=8, color=MUTED)
    # 播放列
    rbox(ax, 0.8, 0.7, 14.3, 1.1, fc=LIGHT, ec=BORDER, rad=0.1)
    tri(ax, 1.15, 1.26, s=0.16, c=TEXT, left=True)
    tri(ax, 1.7, 1.26, s=0.16, c=TEXT)
    ax.add_patch(Rectangle((2.4, 1.15), 9.5, 0.22, fc="#cfd8dc", ec="none"))
    ax.add_patch(Rectangle((2.4, 1.15), 5.2, 0.22, fc=ACCENT, ec="none"))
    ax.add_patch(Circle((7.6, 1.26), 0.16, fc=ACCENT, ec="white", lw=1.2))
    tri(ax, 12.25, 1.26, s=0.15, c=GREEN)
    ax.text(12.5, 1.25, "播放　速度 30fps", fontsize=9, color=TEXT, va="center")
    return _save(fig, "digiinfo.png")


# ---------------------------------------------------------------- record / replay flow
def _flow(ax, steps, y=5.0, colors=None):
    n = len(steps)
    bw, gap = 2.7, 0.9
    total = n * bw + (n - 1) * gap
    x = (16 - total) / 2
    for i, s in enumerate(steps):
        c = (colors or [ACCENT2] * n)[i]
        rbox(ax, x, y - 0.9, bw, 1.8, fc="white", ec=c, lw=2.2, rad=0.14)
        ax.text(x + bw / 2, y, s, fontsize=10.5, ha="center", va="center",
                color=TEXT, fontweight="bold", linespacing=1.4)
        if i < n - 1:
            arrow(ax, x + bw + 0.12, y, x + bw + gap - 0.12, y, c=MUTED)
        x += bw + gap


def fig_record():
    fig, ax = _ax(11, 4.6)
    ax.text(8, 9.0, "監聽錄製 → 自包含 .hidrec → 回放", fontsize=13,
            ha="center", fontweight="bold", color=TEXT)
    _flow(ax, ["監聽\nMonitor", "錄製\nRecord", "匯出\n.hidrec", "載入\n回放"],
          y=5.2, colors=[ACCENT2, ACCENT, GREEN, "#9b59b6"])
    ax.text(8, 2.2, ".hidrec 內含 descriptor + 封包，換台電腦也能完整回放",
            fontsize=9.5, ha="center", color=MUTED)
    return _save(fig, "record.png")


def fig_record_fae():
    fig, ax = _ax(11, 4.6)
    ax.text(8, 9.0, "FAE：監聽 → 錄製 → 匯出（不含回放/載入）", fontsize=13,
            ha="center", fontweight="bold", color=TEXT)
    _flow(ax, ["監聽\nMonitor", "錄製\nRecord", "匯出\n.hidrec"],
          y=5.2, colors=[ACCENT2, ACCENT, GREEN])
    ax.text(8, 2.2, "錄到的檔可交回原廠（工程版）做回放與深入分析",
            fontsize=9.5, ha="center", color=MUTED)
    return _save(fig, "record_fae.png")


# ---------------------------------------------------------------- stress / send
def fig_stress():
    import numpy as np
    fig, ax = _ax(10, 6.4)
    rbox(ax, 0.4, 0.5, 15.2, 9.0, fc="white", ec=BORDER, lw=1.5)
    ax.text(0.8, 9.0, "發送 Send · 壓力測試 Stress", fontsize=13, fontweight="bold", color=TEXT)
    rbox(ax, 0.8, 5.6, 7.0, 2.8, fc=LIGHT, ec=BORDER)
    ax.text(1.1, 7.9, "發送 Output / Feature Report", fontsize=10, fontweight="bold", color=TEXT)
    ax.text(1.1, 7.2, "• 自訂 Report ID + 資料位元組", fontsize=9, color=TEXT)
    ax.text(1.1, 6.65, "• Output 類型可「等待 INT 回應」", fontsize=9, color=TEXT)
    ax.text(1.1, 6.1, "• 即時 log 回傳結果", fontsize=9, color=TEXT)
    rbox(ax, 8.2, 5.6, 7.0, 2.8, fc=LIGHT, ec=BORDER)
    ax.text(8.5, 7.9, "壓力測試", fontsize=10, fontweight="bold", color=TEXT)
    ax.text(8.5, 7.2, "• 連續送命令 + 統計 scan/s", fontsize=9, color=TEXT)
    ax.text(8.5, 6.65, "• 記錄回應延遲", fontsize=9, color=TEXT)
    ax.text(8.5, 6.1, "• 匯出 CSV", fontsize=9, color=TEXT)
    # 折線示意
    rbox(ax, 0.8, 0.9, 14.4, 4.2, fc="#fbfcfd", ec=BORDER)
    ax.text(1.1, 4.7, "scan/s 即時曲線", fontsize=9.5, color=MUTED)
    xs = np.linspace(1.3, 14.8, 120)
    ys = 2.6 + 1.2 * np.sin(np.linspace(0, 9, 120)) + 0.3 * np.random.RandomState(1).randn(120)
    ax.plot(xs, ys, color=ACCENT, lw=1.8)
    return _save(fig, "stress.png")


# ---------------------------------------------------------------- editions compare
def fig_editions():
    fig, ax = _ax(11, 6.6)
    feats = [
        ("監聽 Monitor Data（全裝置 / 單一）", True, True),
        ("Report Descriptor 欄位樹", True, True),
        ("RAW / FF01 欄位", True, True),
        ("監聽畫布（觸控 / 筆視覺化）", True, True),
        ("錄製 + 匯出 .hidrec", True, True),
        ("發送 Send / 壓力測試 Stress", True, False),
        ("回放 Replay（DigiInfo / Differ）", True, False),
        ("載入 .hidrec 回放", True, False),
        ("登入帳號管控", True, False),
    ]
    ax.text(4.6, 9.4, "Engineer\n工程完整版", fontsize=12, ha="center",
            fontweight="bold", color=ACCENT2, linespacing=1.2)
    ax.text(7.7, 9.4, "FAE\n客戶版", fontsize=12, ha="center",
            fontweight="bold", color=ACCENT, linespacing=1.2)
    y = 8.2
    for name, eng, fae in feats:
        ax.text(0.6, y, name, fontsize=10, va="center", color=TEXT)
        check(ax, 4.6, y, GREEN) if eng else cross(ax, 4.6, y, "#bbbbbb")
        check(ax, 7.7, y, GREEN) if fae else cross(ax, 7.7, y, ACCENT)
        ax.plot([0.5, 8.4], [y - 0.33, y - 0.33], color="#eceff1", lw=1)
        y -= 0.72
    ax.text(11.0, 8.2, "FAE 版特點", fontsize=11, fontweight="bold", color=ACCENT)
    for i, t in enumerate(["免帳密、開啟即用", "只看 / 只錄、不回放",
                            "交付客戶現場除錯", "檔案回傳原廠分析"]):
        ax.text(9.0, 7.4 - i * 0.7, "• " + t, fontsize=9.6, color=TEXT)
    return _save(fig, "editions.png")


def fig_agenda():
    fig, ax = _ax(10, 6.4)
    cats = [
        ("1", "裝置 & Report Descriptor 欄位樹", ACCENT2),
        ("2", "監聽模式：全裝置 / 單一裝置", ACCENT2),
        ("3", "Monitor Data 解碼欄位", ACCENT2),
        ("4", "RAW / FF01 廠商欄位", ACCENT2),
        ("5", "監聽畫布：即時觸控 / 筆視覺化", GREEN),
        ("6", "畫布：ID 配色 / Confidence / Eraser", GREEN),
        ("7", "Differ：矩陣差異比對", GREEN),
        ("8", "監聽錄製 → .hidrec 匯出", "#9b59b6"),
        ("9", "回放：載入 .hidrec", "#9b59b6"),
        ("10", "回放 · DigiInfo 手 / 筆雙畫布", "#9b59b6"),
        ("11", "發送 Send / 等待 INT 回應", ACCENT),
        ("12", "壓力測試 Stress", ACCENT),
    ]
    ax.text(0.7, 9.4, "功能總覽", fontsize=16, fontweight="bold", color=TEXT)
    col_x = [0.7, 8.3]
    for i, (num, name, col) in enumerate(cats):
        cx = col_x[i // 6]
        yy = 8.3 - (i % 6) * 1.35
        ax.add_patch(Circle((cx + 0.35, yy), 0.34, fc=col, ec="none"))
        ax.text(cx + 0.35, yy, num, fontsize=12, ha="center", va="center",
                color="white", fontweight="bold")
        ax.text(cx + 0.95, yy, name, fontsize=11.5, va="center", color=TEXT)
    return _save(fig, "agenda.png")


def fig_differ():
    import numpy as np
    fig, ax = _ax(10, 6.4)
    rbox(ax, 0.4, 0.5, 15.2, 9.0, fc="white", ec=BORDER, lw=1.5)
    ax.text(0.8, 9.0, "Differ · 矩陣差異比對", fontsize=13, fontweight="bold", color=TEXT)
    rng = np.random.RandomState(3)
    base = rng.randint(20, 40, (8, 14))
    diff = base.copy()
    diff[3:5, 6:9] += rng.randint(15, 30, (2, 3))
    x0, y0, cw, ch = 0.9, 1.2, 1.02, 0.82
    for r in range(8):
        for c in range(14):
            val = diff[r, c]
            hot = val - base[r, c]
            fc = "#fdecea" if hot > 0 else "white"
            ec = ACCENT if hot > 0 else "#e3e8ec"
            ax.add_patch(Rectangle((x0 + c * cw, y0 + (7 - r) * ch), cw, ch,
                         fc=fc, ec=ec, lw=1.0))
            ax.text(x0 + c * cw + cw / 2, y0 + (7 - r) * ch + ch / 2, str(val),
                    fontsize=7.5, ha="center", va="center",
                    color=ACCENT if hot > 0 else "#607d8b",
                    fontweight="bold" if hot > 0 else "normal")
    ax.text(0.9, 0.8, "兩份資料逐格相減，差異處以紅色標出（找漂移 / 雜訊 / 異常）",
            fontsize=9, color=MUTED)
    return _save(fig, "differ.png")


def fig_send():
    fig, ax = _ax(10, 6.4)
    rbox(ax, 0.4, 0.5, 15.2, 9.0, fc="white", ec=BORDER, lw=1.5)
    ax.text(0.8, 9.0, "發送 Send · Output / Feature Report", fontsize=13,
            fontweight="bold", color=TEXT)
    rbox(ax, 0.8, 6.2, 14.4, 2.2, fc=LIGHT, ec=BORDER)
    ax.text(1.1, 7.9, "Report 類型", fontsize=9.5, color=MUTED)
    for i, (t, on) in enumerate([("Output", True), ("Feature", False), ("Input", False)]):
        rbox(ax, 1.1 + i * 2.2, 6.8, 1.9, 0.7, fc=ACCENT if on else "#eceff1",
             ec="none", rad=0.08)
        ax.text(1.1 + i * 2.2 + 0.95, 7.15, t, fontsize=9.5, ha="center", va="center",
                color="white" if on else "#607d8b",
                fontweight="bold" if on else "normal")
    ax.add_patch(Rectangle((8.2, 6.85), 0.32, 0.32, fc=GREEN, ec="none"))
    check(ax, 8.36, 7.0, "white")
    ax.text(8.7, 7.0, "等待 INT 回應（Output 預設勾選）", fontsize=9.5, va="center", color=TEXT)
    # 資料位元組
    ax.text(1.1, 5.6, "Report ID  0x09", fontsize=10, color=TEXT, family="Consolas")
    for i in range(10):
        rbox(ax, 1.1 + i * 1.35, 4.4, 1.15, 0.8, fc="white", ec=ACCENT2, rad=0.06)
        ax.text(1.1 + i * 1.35 + 0.57, 4.8, f"{(i*17) & 0xFF:02X}", fontsize=10,
                ha="center", va="center", color=TEXT, family="Consolas")
    # log
    rbox(ax, 0.8, 0.9, 14.4, 3.0, fc="#1f2d3d", ec="none")
    logs = ["> Output  ID=0x09  10 bytes", "  [OK] sent",
            "< INT reply  8 bytes:  01 2A 00 ...", "  scan/s = 124"]
    for i, t in enumerate(logs):
        ax.text(1.1, 3.4 - i * 0.62, t, fontsize=9.5,
                color="#8bc34a" if ("OK" in t or "INT" in t) else "#cfd8dc",
                family="Consolas", va="center")
    return _save(fig, "send.png")


def build_figs():
    return {
        "icon": fig_icon(),
        "window": fig_window(),
        "descriptor": fig_descriptor(),
        "digiinfo": fig_digiinfo(),
        "record": fig_record(),
        "record_fae": fig_record_fae(),
        "stress": fig_stress(),
        "editions": fig_editions(),
        "agenda": fig_agenda(),
        "differ": fig_differ(),
        "send": fig_send(),
    }


def placeholder(name, label):
    fig, ax = _ax(10, 6.4)
    ax.add_patch(FancyBboxPatch((0.6, 0.7), 14.8, 8.4,
                 boxstyle="round,pad=0,rounding_size=0.2",
                 fc="#f7f9fb", ec=ACCENT2, lw=2.2, ls=(0, (7, 5))))
    ax.add_patch(FancyBboxPatch((7.0, 5.3), 2.0, 1.5,
                 boxstyle="round,pad=0,rounding_size=0.18", fc="white", ec=MUTED, lw=2))
    ax.add_patch(Circle((8.0, 6.0), 0.45, fc="none", ec=MUTED, lw=2))
    ax.add_patch(Circle((8.0, 6.0), 0.18, fc=MUTED, ec="none"))
    ax.add_patch(Rectangle((7.35, 6.75), 0.55, 0.22, fc=MUTED, ec="none"))
    ax.text(8.0, 4.3, "實機截圖", fontsize=20, ha="center", color=ACCENT2, fontweight="bold")
    ax.text(8.0, 3.5, label, fontsize=13, ha="center", color=TEXT)
    ax.text(8.0, 2.5, f"請置入  slides/shots/{name}.png  後重跑 make_slides.py",
            fontsize=9.5, ha="center", color=MUTED)
    return _save(fig, f"ph_{name}.png")


def shot(name, label, fallback=None):
    """有 slides/shots/<name>.(png/jpg) 就用實機截圖；否則用 fallback 示意圖或佔位圖。"""
    for ext in (".png", ".jpg", ".jpeg", ".PNG", ".JPG"):
        p = os.path.join(SHOTS, name + ext)
        if os.path.exists(p):
            return p
    return fallback or placeholder(name, label)


# ================================================================ PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

JHENG = "Microsoft JhengHei"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def rgb(h):
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def _set_run(r, text, size, color, bold=False, font=JHENG):
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = rgb(color); r.font.name = font


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background(); shape.shadow.inherit = False


def add_fitted(slide, path, left, top, box_w, box_h, border=True):
    """把圖等比例縮放置中放入 (box_w × box_h)，截圖長寬不一也不變形。"""
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = box_w; h = int(w / ar)
    if h > box_h:
        h = box_h; w = int(h * ar)
    px = left + (box_w - w) // 2
    py = top + (box_h - h) // 2
    pic = slide.shapes.add_picture(path, px, py, width=w, height=h)
    if border:
        pic.line.color.rgb = rgb(BORDER); pic.line.width = Pt(1)
    return pic


def cover(prs, title, subtitle, tagline, icon, accent):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H); _fill(bg, PRIMARY)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), EMU_W, Inches(0.12)); _fill(stripe, accent)
    if icon and os.path.exists(icon):
        s.shapes.add_picture(icon, Inches(0.9), Inches(0.8), height=Inches(1.5))
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; _set_run(p.add_run(), title, 40, "#ffffff", bold=True)
    p2 = tf.add_paragraph(); _set_run(p2.add_run(), subtitle, 22, accent, bold=True); p2.space_before = Pt(8)
    p3 = tf.add_paragraph(); _set_run(p3.add_run(), tagline, 15, "#b0bec5"); p3.space_before = Pt(14)
    return s


def section(prs, text, accent):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H); _fill(bg, PRIMARY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.2), Inches(0.18), Inches(1.1)); _fill(bar, accent)
    tb = s.shapes.add_textbox(Inches(1.3), Inches(3.1), Inches(11), Inches(1.4))
    _set_run(tb.text_frame.paragraphs[0].add_run(), text, 32, "#ffffff", bold=True)
    return s


def content(prs, title, bullets, image=None, accent=ACCENT, img_left=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(1.05)); _fill(bar, PRIMARY)
    st = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), EMU_W, Inches(0.08)); _fill(st, accent)
    tbar = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85))
    tbar.text_frame.word_wrap = True
    pp = tbar.text_frame.paragraphs[0]; pp.anchor = None
    _set_run(pp.add_run(), title, 26, "#ffffff", bold=True)

    txt_w = Inches(5.2) if image else Inches(12.3)
    tx = Inches(7.4) if (image and img_left) else Inches(0.6)
    ix = Inches(0.5) if (image and img_left) else Inches(6.0)
    tb = s.shapes.add_textbox(tx, Inches(1.45), txt_w, Inches(5.7))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        lvl = 0
        if isinstance(b, tuple):
            b, lvl = b
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        if lvl == 0:
            run = p.add_run(); _set_run(run, "▍", 14, accent, bold=True)
            _set_run(p.add_run(), " " + b, 15.5, TEXT, bold=True)
        else:
            _set_run(p.add_run(), "   • " + b, 13.5, "#37474f")
            p.space_after = Pt(4)
    if image and os.path.exists(image):
        add_fitted(s, image, ix, Inches(1.5), Inches(6.9), Inches(5.5))
    return s


def closing(prs, accent, lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H); _fill(bg, PRIMARY)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.3), EMU_W, Inches(0.12)); _fill(stripe, accent)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.4))
    tf = tb.text_frame
    _set_run(tf.paragraphs[0].add_run(), lines[0], 34, "#ffffff", bold=True)
    for ln in lines[1:]:
        p = tf.add_paragraph(); _set_run(p.add_run(), ln, 16, "#b0bec5"); p.space_before = Pt(10)
    return s


def deck_full(figs):
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    cover(prs, "RE024 Touch Inspector",
          "完整功能版（Engineer Edition）",
          "HID 觸控 / 手寫筆 裝置監聽・解碼・視覺化・回放工具", figs["icon"], ACCENT2)

    # ── 功能總覽（先說明有哪些功能）──
    content(prs, "功能總覽",
            ["本工具一站式涵蓋觸控 / 手寫筆裝置的：",
             ("① 裝置與 Descriptor 解析", 1),
             ("② 即時監聽（表格 + 畫布）", 1),
             ("③ 錄製與回放（含 DigiInfo log）", 1),
             ("④ 主動發送與壓力測試", 1),
             "以下逐一說明各功能　→"],
            figs["agenda"], ACCENT2)

    # ── 逐一功能（每頁一個）──
    content(prs, "① 裝置 & Report Descriptor 欄位樹",
            ["選擇要分析的 HID 裝置（USB / I2C-HID）",
             "解析 Report Descriptor → 樹狀欄位",
             ("Collection / Usage 階層一目了然", 1),
             ("每欄位的 Usage、位元位置、邏輯範圍", 1),
             ("X/Y、ContactID、TipSwitch、Confidence、Pressure、Tilt…", 1)],
            shot("01_descriptor", "裝置清單 + Report Descriptor 欄位樹", figs["descriptor"]),
            ACCENT2, img_left=False)

    content(prs, "② 監聽模式：全裝置 / 單一裝置",
            ["全裝置 digitizer 自動模式（預設）",
             ("不選裝置，自動解碼系統上所有 digitizer", 1),
             ("每裝置獨立 scan/s 統計", 1),
             "或鎖定單一裝置深入觀察",
             ("只處理選定裝置封包，不受其他裝置干擾", 1)],
            shot("02_monitor_modes", "監聽分頁 — 全裝置 / 單一裝置切換", figs["window"]),
            ACCENT2)

    content(prs, "③ Monitor Data 解碼欄位",
            ["即時逐封包解碼成表格",
             ("X / Y、CenterX / CenterY", 1),
             ("ContactID、Status（Tip/InRange/Eraser…）", 1),
             ("Confidence、Pressure、Tilt、Azimuth", 1),
             "可匯出 CSV / Excel"],
            shot("03_monitor_data", "Monitor Data 解碼欄位表格", figs["window"]),
            ACCENT, img_left=False)

    content(prs, "④ RAW / FF01 廠商欄位",
            ["RAW：顯示每個 byte 原始值",
             ("逐位元組對照、debug 底層資料", 1),
             "FF01：廠商自訂（Vendor）欄位",
             ("解出廠商 debug / 校正資訊", 1)],
            shot("04_raw_ff01", "RAW / FF01 欄位顯示", figs["window"]),
            ACCENT)

    content(prs, "⑤ 監聽畫布：即時觸控 / 筆視覺化",
            ["畫布與 Monitor Data 各佔一半畫面",
             "每裝置一格 grid、即時描繪接點與軌跡",
             ("觸控點移動軌跡、手寫筆筆跡", 1),
             "確認 tip-off / inrange-off 才清除接點",
             ("多裝置交錯也不誤清軌跡", 1)],
            shot("05_canvas", "監聽畫布 — 即時觸控 / 筆", figs["window"]),
            GREEN, img_left=False)

    content(prs, "⑥ 畫布：ID 配色 / Confidence / Eraser",
            ["每個接點 ID 不同顏色，易追蹤",
             ("Confidence=0（低信心 / palm）→ 接點放大加粗", 1),
             ("手寫筆 tip 或 eraser 接觸時都出線", 1),
             "快速看出 palm rejection 與筆觸狀態"],
            shot("06_canvas_detail", "畫布 — ID 配色 / Confidence 放大", figs["window"]),
            GREEN)

    content(prs, "⑦ Differ：矩陣差異比對",
            ["兩份資料逐格相減",
             ("差異處以紅色 heatmap 標出", 1),
             "找漂移 / 雜訊 / 異常點",
             ("支援大矩陣（如 46×42）顯示數值", 1)],
            shot("07_differ", "Differ 矩陣差異比對", figs["differ"]),
            GREEN, img_left=False)

    content(prs, "⑧ 監聽錄製 → .hidrec 匯出",
            ["一鍵錄製監聽資料",
             "匯出自包含 .hidrec 檔",
             ("內含 descriptor + 封包", 1),
             ("換台電腦也能完整回放", 1)],
            shot("08_record", "監聽錄製 / 匯出 .hidrec", figs["record"]),
            "#9b59b6")

    content(prs, "⑨ 回放：載入 .hidrec",
            ["載入先前錄製的 .hidrec",
             "重現當時監聽畫面與資料",
             ("切換分頁自動暫停，效能友善", 1),
             "可反覆檢視、逐幀分析"],
            shot("09_replay_load", "載入 .hidrec 回放", figs["record"]),
            "#9b59b6", img_left=False)

    content(prs, "⑩ 回放 · DigiInfo（手 / 筆雙畫布）",
            ["載入 DigiInfo 觸控 log（XML）",
             "手（觸控）/ 筆 分兩張畫布、座標軸獨立",
             ("座標範圍取自 log 表頭、等比例不變形", 1),
             ("筆只在 tip / eraser / 有壓力 時連線", 1),
             "逐幀播放 + 軌跡 + 資料表",
             ("表格含 tip / invert / eraser / inrange / pressure", 1)],
            shot("10_digiinfo", "DigiInfo 手 / 筆雙畫布 + 播放", figs["digiinfo"]),
            "#9b59b6")

    content(prs, "⑪ 發送 Send / 等待 INT 回應",
            ["發送 Output / Feature Report",
             ("自訂 Report ID 與資料位元組", 1),
             ("Output 類型可「等待 INT 回應」（預設勾選）", 1),
             "即時 log 顯示送出與回應結果"],
            shot("11_send", "發送 Send + 等待 INT 回應", figs["send"]),
            ACCENT, img_left=False)

    content(prs, "⑫ 壓力測試 Stress",
            ["連續送命令壓力測試",
             ("即時統計 scan/s 與回應延遲", 1),
             ("曲線即時呈現", 1),
             "結果可匯出 CSV"],
            shot("12_stress", "壓力測試 scan/s 曲線", figs["stress"]),
            ACCENT)

    # ── 版本比較 + 結尾 ──
    content(prs, "版本比較：Engineer / FAE",
            ["Engineer：完整功能 + 登入帳號管控",
             "FAE：客戶 / 現場用的閹割版",
             ("免帳密、只監聽 / 只錄製", 1),
             ("隱藏發送 / 壓測 / 回放", 1),
             "同一套程式，建置時切換版本"],
            figs["editions"], ACCENT, img_left=False)

    closing(prs, ACCENT2,
            ["RE024 Touch Inspector",
             "完整功能版 — 觸控 / 手寫筆 裝置的工程利器",
             "監聽 · 解碼 · 視覺化 · 錄製 · 回放"])
    return prs


def deck_fae(figs):
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    cover(prs, "RE024 Touch Inspector (FAE)",
          "客戶版 / 現場版（FAE Edition）",
          "免登入・觸控與手寫筆監聽・錄製交付原廠分析", figs["icon"], ACCENT)

    content(prs, "FAE 版定位",
            ["為客戶 / FAE 現場除錯而生的精簡版",
             ("開啟即用，免帳號密碼", 1),
             "能做：監聽、看欄位、看畫布、錄製匯出",
             "不含：發送、壓測、回放、載入錄製",
             ("把錄到的 .hidrec 回傳原廠（工程版）深入分析", 1)],
            figs["editions"], ACCENT)

    content(prs, "監聽 Monitor Data",
            ["全裝置 digitizer 自動模式 / 單一裝置",
             ("不選裝置也能自動解碼所有 digitizer", 1),
             "解碼欄位 + RAW / FF01",
             ("X/Y、ContactID、Status、Confidence", 1),
             ("每裝置獨立 scan/s", 1)],
            shot("03_monitor_data", "Monitor Data 解碼欄位表格", figs["window"]), ACCENT2)

    content(prs, "Report Descriptor 欄位樹",
            ["解析裝置 Report Descriptor",
             ("樹狀看 Collection / Usage / 範圍", 1),
             "快速確認觸控 / 手寫筆欄位定義",
             ("X/Y、ContactID、TipSwitch、Confidence、Pressure…", 1)],
            shot("01_descriptor", "裝置清單 + Report Descriptor 欄位樹", figs["descriptor"]),
            ACCENT2, img_left=False)

    content(prs, "監聽畫布（即時視覺化）",
            ["每裝置一格、與表格各佔一半畫面",
             "每個接點 ID 不同顏色 + 軌跡",
             ("手寫筆 tip / eraser 接觸時出線", 1),
             ("Confidence=0 接點放大，方便看 palm", 1)],
            shot("05_canvas", "監聽畫布 — 即時觸控 / 筆", figs["window"]), GREEN, img_left=False)

    content(prs, "錄製與交付",
            ["一鍵錄製現場監聽資料",
             "匯出自包含 .hidrec 檔",
             ("含 descriptor + 封包", 1),
             "回傳原廠以工程版回放、深入分析",
             ("FAE 版本身不提供回放 / 載入", 1)],
            shot("08_record", "監聽錄製 / 匯出 .hidrec", figs["record_fae"]), ACCENT)

    closing(prs, ACCENT,
            ["RE024 Touch Inspector (FAE)",
             "現場監聽 · 即時視覺化 · 錄製交付",
             "簡單、免登入、給客戶與 FAE 使用"])
    return prs


def main():
    figs = build_figs()
    deck_full(figs).save(os.path.join(OUT, "RE024_Touch_Inspector_完整功能版.pptx"))
    deck_fae(figs).save(os.path.join(OUT, "RE024_Touch_Inspector_FAE客戶版.pptx"))
    print("OK ->", OUT)


if __name__ == "__main__":
    main()
